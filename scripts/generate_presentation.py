"""Generate updated project presentation as PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Color palette (matching existing PDF)
# ---------------------------------------------------------------------------
DARK_BLUE   = RGBColor(0x1E, 0x3A, 0x5F)
GREEN       = RGBColor(0x2D, 0x6A, 0x3F)
TEAL        = RGBColor(0x00, 0x6D, 0x77)
PURPLE      = RGBColor(0x5B, 0x21, 0x86)
RED_ACCENT  = RGBColor(0xB9, 0x1C, 0x1C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF6, 0xF9)
MID_GRAY    = RGBColor(0x66, 0x70, 0x85)
BORDER      = RGBColor(0xD0, 0xD7, 0xE3)
DARK_TEXT   = RGBColor(0x1A, 0x20, 0x2C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, x, y, w, h, text, size=14, bold=False, color=DARK_TEXT,
            align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, Inches(1.05), fill=DARK_BLUE)
    textbox(slide, Inches(0.4), Inches(0.1), Inches(11), Inches(0.55),
            title, size=24, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, Inches(0.4), Inches(0.62), Inches(11), Inches(0.35),
                subtitle, size=12, color=RGBColor(0xA0, 0xB4, 0xCC))


def card(slide, x, y, w, h, title=None, title_color=GREEN, body_lines=None,
         bg=LIGHT_GRAY, border=BORDER):
    rect(slide, x, y, w, h, fill=bg, line=border)
    cy = y + Inches(0.15)
    if title:
        textbox(slide, x + Inches(0.2), cy, w - Inches(0.3), Inches(0.3),
                title, size=10, bold=True, color=title_color)
        cy += Inches(0.32)
        rect(slide, x + Inches(0.2), cy, w - Inches(0.4), Pt(1), fill=BORDER)
        cy += Inches(0.1)
    if body_lines:
        for line in body_lines:
            if line == "---":
                cy += Inches(0.08)
                continue
            size = 9 if not line.startswith("*") else 10
            bold = line.startswith("*")
            text = line.lstrip("*")
            textbox(slide, x + Inches(0.2), cy, w - Inches(0.35), Inches(0.28),
                    text, size=size, bold=bold, color=DARK_TEXT if bold else MID_GRAY)
            cy += Inches(0.25)


def numbered_box(slide, x, y, size_px, number, bg_color):
    rect(slide, x, y, Inches(0.42), Inches(0.42), fill=bg_color)
    textbox(slide, x, y, Inches(0.42), Inches(0.42),
            str(number), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs):
    sl = blank_slide(prs)
    # Left dark panel
    rect(sl, 0, 0, Inches(5), SLIDE_H, fill=DARK_BLUE)
    # Right white area
    rect(sl, Inches(5), 0, Inches(8.33), SLIDE_H, fill=WHITE)
    # CO2 badge top-right
    rect(sl, Inches(11.5), Inches(0.3), Inches(1.5), Inches(1.2), fill=DARK_BLUE)
    textbox(sl, Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.7),
            "CO₂", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, Inches(11.5), Inches(1.0), Inches(1.5), Inches(0.35),
            "Budget Monitor", size=9, color=WHITE, align=PP_ALIGN.CENTER)
    # Left panel text
    textbox(sl, Inches(0.4), Inches(1.5), Inches(4.2), Inches(0.45),
            "Visualisierungsprojekt", size=14, color=RGBColor(0xA0, 0xB4, 0xCC))
    # Title box on right side
    rect(sl, Inches(5.3), Inches(1.6), Inches(7.5), Inches(2.4), fill=WHITE, line=BORDER)
    textbox(sl, Inches(5.5), Inches(1.7), Inches(7.1), Inches(1.6),
            "Business Travel\nEmissions Dashboard", size=32, bold=True, color=DARK_BLUE)
    textbox(sl, Inches(5.5), Inches(3.3), Inches(7.1), Inches(0.35),
            "Einreichung Entwurf Präsentation  ·  18. Mai 2026", size=11, color=MID_GRAY)
    textbox(sl, Inches(5.5), Inches(3.7), Inches(7.1), Inches(0.35),
            "Fiton Sallhaj  ·  Christopher Zickler  ·  Villavarayasingam Gopigan",
            size=11, color=MID_GRAY)


def slide_thema(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Thema & Hauptziel")

    # Left card: THEMA
    rect(sl, Inches(0.3), Inches(1.2), Inches(6.0), Inches(5.9), fill=WHITE, line=BORDER)
    textbox(sl, Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.3),
            "THEMA", size=10, bold=True, color=GREEN)
    textbox(sl, Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.8),
            "CO₂-Emissionen aus Geschäftsreisen —\noperatives Monitoring-Dashboard",
            size=14, bold=True, color=DARK_BLUE)
    rect(sl, Inches(0.5), Inches(2.55), Inches(5.6), Pt(1), fill=BORDER)

    details = [
        ("Produktform",   "Interaktives Streamlit-Dashboard (Desktop)"),
        ("Datenbasis",    "Alle Geschäftsreisen 2017–2026\n(Flug, Bahn, Bus, Mietwagen)"),
        ("Update-Zyklus", "Monatlich (kein Echtzeit-Tracking)"),
    ]
    cy = Inches(2.7)
    for label, val in details:
        textbox(sl, Inches(0.5), cy, Inches(5.5), Inches(0.25),
                label, size=10, bold=True, color=DARK_TEXT)
        cy += Inches(0.25)
        textbox(sl, Inches(0.5), cy, Inches(5.5), Inches(0.45),
                val, size=10, color=MID_GRAY)
        cy += Inches(0.5)

    textbox(sl, Inches(0.5), Inches(6.5), Inches(5.5), Inches(0.35),
            "Sekundär: ESG-Reporting (Swiss Art. 964a OR / EU CSRD)",
            size=9, color=MID_GRAY)

    # Right card: HAUPTZIEL
    rect(sl, Inches(6.7), Inches(1.2), Inches(6.3), Inches(5.9), fill=WHITE, line=BORDER)
    textbox(sl, Inches(6.9), Inches(1.35), Inches(5.8), Inches(0.3),
            "HAUPTZIEL", size=10, bold=True, color=GREEN)
    textbox(sl, Inches(6.9), Inches(1.7), Inches(5.8), Inches(0.45),
            "Team Heads in < 1 Minute befähigen zu sehen:",
            size=13, bold=True, color=DARK_BLUE)

    goals = [
        (GREEN,      "?",  "Status",   "Wie viel Budget haben wir noch?"),
        (TEAL,       "↓",  "Ursache",  "Welche Trips treiben unsere Emissionen?"),
        (DARK_BLUE,  "✓",  "Handlung", "Wo können wir konkret sparen?"),
    ]
    cy = Inches(2.4)
    for color, icon, label, desc in goals:
        rect(sl, Inches(6.9), cy, Inches(0.5), Inches(0.5), fill=color)
        textbox(sl, Inches(6.9), cy, Inches(0.5), Inches(0.5),
                icon, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(sl, Inches(7.5), cy, Inches(5.2), Inches(0.28),
                label, size=12, bold=True, color=color)
        textbox(sl, Inches(7.5), cy + Inches(0.28), Inches(5.2), Inches(0.28),
                desc, size=10, color=MID_GRAY)
        cy += Inches(0.85)


def slide_zielgruppe(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Zielgruppe")

    # Left big card: PRIMÄRNUTZER
    rect(sl, Inches(0.3), Inches(1.2), Inches(8.2), Inches(5.9), fill=WHITE, line=BORDER)
    textbox(sl, Inches(0.5), Inches(1.35), Inches(7.8), Inches(0.3),
            "PRIMÄRNUTZER", size=10, bold=True, color=GREEN)

    # Persona avatar box
    rect(sl, Inches(0.5), Inches(1.75), Inches(1.4), Inches(1.4), fill=DARK_BLUE)
    textbox(sl, Inches(0.5), Inches(1.9), Inches(1.4), Inches(1.0),
            "👤", size=28, align=PP_ALIGN.CENTER)

    textbox(sl, Inches(2.1), Inches(1.75), Inches(6.2), Inches(0.45),
            "Team Head / BU-Manager", size=16, bold=True, color=DARK_BLUE)
    textbox(sl, Inches(2.1), Inches(2.2), Inches(6.2), Inches(0.35),
            "z. B. Sarah, Sales BU – 15-köpfiges Team, jährliches CO₂-Budget",
            size=10, color=MID_GRAY)
    rect(sl, Inches(0.5), Inches(3.25), Inches(7.8), Pt(1), fill=BORDER)

    rows = [
        ("Frequenz",       "1–2× pro Monat + ad hoc vor Reisegenehmigungen"),
        ("Datenaffinität", "Gering – braucht klare, handlungsorientierte Aussagen"),
        ("Hauptschmerz",   "Kein Überblick bis zum Quartalsbericht – zu spät zum Handeln"),
        ("Ziel",           "In 2 Min. wissen: Sind wir im Budget? Wer treibt es?"),
    ]
    cy = Inches(3.45)
    for label, val in rows:
        textbox(sl, Inches(0.5), cy, Inches(2.0), Inches(0.28),
                label, size=10, bold=True, color=GREEN)
        textbox(sl, Inches(2.6), cy, Inches(5.7), Inches(0.28),
                val, size=10, color=DARK_TEXT)
        cy += Inches(0.45)

    # Right: Sekundär + Nicht im Scope
    rect(sl, Inches(8.8), Inches(1.2), Inches(4.2), Inches(3.1), fill=WHITE, line=BORDER)
    textbox(sl, Inches(9.0), Inches(1.35), Inches(3.8), Inches(0.3),
            "SEKUNDÄRNUTZER", size=10, bold=True, color=TEAL)

    sec = [
        ("CSO (Dr. Müller)",
         "Aggregierte Monatsberichte\nfür ESG-Compliance (GRI/ESRS)"),
        ("Finance (Marc)",
         "CSV-Export für Audit-Trail\n& Budgetvalidierung"),
    ]
    cy = Inches(1.75)
    for name, desc in sec:
        textbox(sl, Inches(9.0), cy, Inches(3.8), Inches(0.28),
                name, size=11, bold=True, color=DARK_BLUE)
        textbox(sl, Inches(9.0), cy + Inches(0.28), Inches(3.8), Inches(0.45),
                desc, size=9, color=MID_GRAY)
        cy += Inches(0.9)

    rect(sl, Inches(8.8), Inches(4.55), Inches(4.2), Inches(2.55),
         fill=RGBColor(0xFF, 0xF1, 0xF0), line=RGBColor(0xFC, 0xC, 0xC))
    textbox(sl, Inches(9.0), Inches(4.7), Inches(3.8), Inches(0.3),
            "NICHT IM SCOPE", size=10, bold=True, color=RED_ACCENT)
    textbox(sl, Inches(9.0), Inches(5.1), Inches(3.8), Inches(1.6),
            "Mitarbeitende (Datenquelle, anonymisiert)\nBreite Öffentlichkeit / Externe\nTop-Management / Executive Reporting",
            size=10, color=DARK_TEXT)


def slide_dashboard(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Dashboard-Struktur", "3 Views · Context → Insight → Action")

    views = [
        (GREEN,      "BU-Leitung\n(Primär)",
         ["Budget-Status KPIs", "Actual vs. Budget Chart", "Emissionstreiber", "Rail-First Chancen"]),
        (TEAL,       "CSO Report\n(Sekundär)",
         ["BU-Vergleich aggregiert", "Trend 2017–2026", "Budget-Übersicht alle BUs"]),
        (MID_GRAY,   "Finance Report\n(Sekundär)",
         ["Kosten-Emissions-Tabelle", "High-cost / High-emission Trips", "CSV-Export"]),
    ]

    cx = Inches(0.3)
    for color, title, items in views:
        rect(sl, cx, Inches(1.2), Inches(4.1), Inches(5.9), fill=WHITE, line=BORDER)
        rect(sl, cx, Inches(1.2), Inches(4.1), Inches(0.65), fill=color)
        textbox(sl, cx + Inches(0.15), Inches(1.25), Inches(3.8), Inches(0.55),
                title, size=13, bold=True, color=WHITE)
        cy = Inches(2.05)
        for item in items:
            rect(sl, cx + Inches(0.2), cy + Inches(0.06), Inches(0.15), Inches(0.15), fill=color)
            textbox(sl, cx + Inches(0.45), cy, Inches(3.5), Inches(0.32),
                    item, size=11, color=DARK_TEXT)
            cy += Inches(0.42)
        cx += Inches(4.45)

    # Filter overview at bottom
    rect(sl, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.3), fill=LIGHT_GRAY, line=BORDER)
    textbox(sl, Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.25),
            "Filter (Sidebar):  Jahr  ·  Monat  ·  Business Unit  ·  Subunit  ·  CO₂e-Szenario (RFI 2.0 / 2.7)  ·  Transportmodus  ·  Reisezweck  ·  Haul",
            size=9, color=MID_GRAY)


def slide_visualisierungen(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Visualisierungstypen")

    charts = [
        (GREEN,     "Grouped Bar Chart",
         "Budget vs. Ist-Emissionen",
         "Sofortiger Längenvergleich statt Winkelinterpretation (besser als Gauge)"),
        (TEAL,      "Bar Charts",
         "Emissionstreiber (3×)",
         "Nach Transportmodus, Reisezweck und Haul – zeigt direkt, was zu tun ist"),
        (DARK_BLUE, "Timeline",
         "Emissionsverlauf 2017–2026",
         "Saisonale Muster und Langzeittrend sichtbar machen"),
        (PURPLE,    "Tabelle + KPIs",
         "Rail-First Chancen",
         "Flüge mit Zugalternative nach CO₂ gerankt; Kosten und km als Kontext"),
    ]

    positions = [(Inches(0.3), Inches(1.2)), (Inches(6.8), Inches(1.2)),
                 (Inches(0.3), Inches(4.4)), (Inches(6.8), Inches(4.4))]

    for (x, y), (color, ctype, title, desc) in zip(positions, charts):
        rect(sl, x, y, Inches(6.2), Inches(2.9), fill=WHITE, line=BORDER)
        rect(sl, x, y, Inches(0.12), Inches(2.9), fill=color)
        textbox(sl, x + Inches(0.25), y + Inches(0.15), Inches(5.7), Inches(0.3),
                ctype, size=13, bold=True, color=color)
        textbox(sl, x + Inches(0.25), y + Inches(0.5), Inches(5.7), Inches(0.3),
                title, size=11, bold=True, color=DARK_TEXT)
        rect(sl, x + Inches(0.25), y + Inches(0.85), Inches(5.7), Pt(1), fill=BORDER)
        textbox(sl, x + Inches(0.25), y + Inches(0.98), Inches(5.7), Inches(1.7),
                desc, size=10, color=MID_GRAY)


def slide_feedback(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Aspekte, zu denen wir Feedback erwarten")

    items = [
        (GREEN,      "1", "Zielgruppe & Workflow",
         "Deckt das Dashboard die typischen Fragen von Team Heads ab?\nIst der «Monday Morning»-Check realistisch in < 2 Min.?"),
        (TEAL,       "2", "Narrative Wirksamkeit",
         "Ist der Context→Insight→Action-Fluss intuitiv?\nSind die 3 Views (BU-Leitung / CSO / Finance) klar getrennt?"),
        (DARK_BLUE,  "3", "Visualisierungen",
         "Grouped Bar Chart für Budget vs. Actual — intuitiv genug?\nSind 3 separate Treiber-Charts besser als ein Sankey?"),
        (PURPLE,     "4", "Unsicherheit (RFI 2.0 vs. 2.7)",
         "Ist die Dual-Szenario-Darstellung verständlich, ohne zu verwirren?"),
        (RED_ACCENT, "5", "Fehlende Features",
         "Brauchen Manager Export-Funktionen (CSV/PDF)?\nMobile-Optimierung nötig?"),
    ]

    positions = [
        (Inches(0.3),  Inches(1.2)),
        (Inches(4.65), Inches(1.2)),
        (Inches(9.0),  Inches(1.2)),
        (Inches(0.3),  Inches(4.35)),
        (Inches(4.65), Inches(4.35)),
    ]

    for (x, y), (color, num, title, body) in zip(positions, items):
        rect(sl, x, y, Inches(4.1), Inches(2.85), fill=WHITE, line=BORDER)
        rect(sl, x, y, Inches(4.1), Inches(0.55), fill=color)
        textbox(sl, x + Inches(0.15), y + Inches(0.08), Inches(0.4), Inches(0.4),
                num, size=16, bold=True, color=WHITE)
        textbox(sl, x + Inches(0.6), y + Inches(0.1), Inches(3.3), Inches(0.38),
                title, size=12, bold=True, color=WHITE)
        textbox(sl, x + Inches(0.2), y + Inches(0.7), Inches(3.7), Inches(2.0),
                body, size=10, color=DARK_TEXT)


def slide_sarah(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Beispiel: Sarah's Monday Morning",
               "9:10 Uhr — Sarah öffnet das Dashboard. Ziel: Budget-Status in 2 Minuten verstehen.")

    steps = [
        (GREEN,     "Landing",    "42 / 50 t CO₂\n(84%)",
         "Große Kennzahl sofort sichtbar · Farbe: Grün = Im Budget"),
        (TEAL,      "Drill-down", "Langstrecken\n= Top-Treiber",
         "Emissionstreiber-Tab: Transport, Zweck, Haul · Balkendiagramm"),
        (DARK_BLUE, "Action",     "Rail-First\nChancen",
         "Routen mit Zugalternative nach CO₂ gerankt · Sarah sieht Top-Route"),
        (GREEN,     "Ergebnis",   "Entscheidung\nin 2 Min.",
         "Sarah priorisiert essenzielle Reisen · nächste Buchung: Zug"),
    ]

    cx = Inches(0.3)
    for i, (color, stage, big, small) in enumerate(steps):
        rect(sl, cx, Inches(1.2), Inches(3.1), Inches(5.9), fill=WHITE, line=BORDER)
        rect(sl, cx, Inches(1.2), Inches(3.1), Inches(0.5), fill=color)
        textbox(sl, cx + Inches(0.15), Inches(1.25), Inches(2.8), Inches(0.38),
                stage, size=12, bold=True, color=WHITE)
        textbox(sl, cx + Inches(0.15), Inches(1.85), Inches(2.8), Inches(1.1),
                big, size=20, bold=True, color=color)
        rect(sl, cx + Inches(0.15), Inches(3.1), Inches(2.7), Pt(1), fill=BORDER)
        textbox(sl, cx + Inches(0.15), Inches(3.25), Inches(2.8), Inches(3.5),
                small, size=10, color=MID_GRAY)
        if i < 3:
            textbox(sl, cx + Inches(3.15), Inches(3.5), Inches(0.25), Inches(0.4),
                    "→", size=18, color=MID_GRAY, align=PP_ALIGN.CENTER)
        cx += Inches(3.35)


def slide_next_steps(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Nächste Schritte & Kontakt")

    # Left: next steps
    rect(sl, Inches(0.3), Inches(1.2), Inches(7.8), Inches(5.9), fill=WHITE, line=BORDER)
    textbox(sl, Inches(0.5), Inches(1.35), Inches(7.4), Inches(0.3),
            "NÄCHSTE SCHRITTE", size=10, bold=True, color=GREEN)

    steps = [
        "MVP-Testing mit 5 Team Heads (Usability-Tests)",
        "Integration von Export-Funktionen (CSV / PDF)",
        "Evaluation und Iteration basierend auf Feedback",
        "Finalisierung & Abgabe bis 1. Juni 2026",
    ]
    cy = Inches(1.85)
    for i, step in enumerate(steps, 1):
        rect(sl, Inches(0.5), cy, Inches(0.45), Inches(0.45), fill=GREEN)
        textbox(sl, Inches(0.5), cy, Inches(0.45), Inches(0.45),
                str(i), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(sl, Inches(1.1), cy + Inches(0.07), Inches(6.8), Inches(0.32),
                step, size=12, color=DARK_TEXT)
        cy += Inches(0.9)

    # Right: team
    rect(sl, Inches(8.4), Inches(1.2), Inches(4.6), Inches(5.9), fill=WHITE, line=BORDER)
    textbox(sl, Inches(8.6), Inches(1.35), Inches(4.2), Inches(0.3),
            "TEAM", size=10, bold=True, color=GREEN)

    members = [
        ("Fiton Sallhaj",              "Project Lead & Storytelling"),
        ("Christopher Zickler",        "Data Engineering & CO₂-Berechnungen"),
        ("Villavarayasingam Gopigan",   "Visualization & UI/UX Design"),
    ]
    cy = Inches(1.85)
    for name, role in members:
        rect(sl, Inches(8.6), cy, Inches(4.0), Inches(1.2), fill=LIGHT_GRAY, line=BORDER)
        textbox(sl, Inches(8.8), cy + Inches(0.15), Inches(3.6), Inches(0.38),
                name, size=12, bold=True, color=DARK_BLUE)
        textbox(sl, Inches(8.8), cy + Inches(0.55), Inches(3.6), Inches(0.45),
                role, size=10, color=MID_GRAY)
        cy += Inches(1.45)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_thema(prs)
    slide_zielgruppe(prs)
    slide_dashboard(prs)
    slide_visualisierungen(prs)
    slide_feedback(prs)
    slide_sarah(prs)
    slide_next_steps(prs)

    out = "docs/presentation_visualization_project.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
